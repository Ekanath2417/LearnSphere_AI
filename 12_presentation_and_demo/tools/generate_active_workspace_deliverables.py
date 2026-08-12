from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "12_presentation_and_demo" / "deliverables"
SCREENSHOT = ROOT / "12_presentation_and_demo" / "assets" / "active-workspace-home.png"
OUTPUT.mkdir(parents=True, exist_ok=True)

TITLE = "LearnSphere AI Active Workspace Booklet"
SUBTITLE = "MindForgeAI Internship 2026 | Chatake Innoworks Pvt. Ltd."
TEAM = "LearnSphere AI Team: Vaishnavi Mali (2026IT012, Team Lead), Nagweni Kumbhar (2026IT011), Aishwarya Ekanath (2026IT010)"

SECTIONS = [
    ("1. Purpose", "The active workspace at D:\\LearnSphere_AI_ActiveWorkspace is the approved place for future changes. It is separate from the baseline project folder and has its own local server on port 5051."),
    ("2. Start and Stop Localhost", "Run .\\04_active_workspace\\RUN_LOCALHOST.ps1 from the active workspace, then open http://localhost:5051. Stop the server with Ctrl+C or run .\\04_active_workspace\\STOP_LOCALHOST.ps1."),
    ("3. Page-Level Code Structure", "The application shell is 06_code/frontend/index.html. Interactive behavior is in assets/app.js and shared styling is in assets/app.css. Understandable information pages are stored separately in 06_code/frontend/pages: overview.html, planner.html, library.html, practice.html, and team.html. Their shared rules are in pages.css."),
    ("4. Change Map", "Use 04_active_workspace/CHANGE_MAP.md before editing. It maps each desired change to its frontend file, backend route, database concern, deployment setting, and verification method."),
    ("5. Integration", "Locally, the browser calls the Flask API through /api. Flask routes work with SQLite. On GitHub Pages, assets/runtime-config.js points to the Render API. Render must allow https://ekanath2417.github.io through CORS_ORIGINS."),
    ("6. Git Commands", "Use git status to see work, git pull origin master to receive approved changes, git switch -c feature/name to create a branch, git add and git commit to record work, and git push -u origin feature/name to send it to GitHub. Always inspect git diff before committing."),
    ("7. Quality Rules", "Run python -m unittest discover -s 08_testing before committing. Never commit .env files, passwords, student uploads, or databases. Keep learning insights supportive and never describe marks as certain or diagnostic."),
    ("8. Deployment", "GitHub Pages publishes the static website. Render hosts the Flask API. Confirm Pages workflow success, Render /api/health response, CORS origin, and runtime API URL before testing hosted sign-in."),
    ("9. Troubleshooting", "If a change is not visible, confirm that localhost:5051 is open, refresh with Ctrl+F5, inspect the terminal log, confirm the file was edited in the active workspace, and consult the change map. If hosted sign-in fails, check Render health, CORS, and runtime-config.js."),
]


def add_docx():
    doc = Document()
    doc.add_heading(TITLE, 0)
    doc.add_paragraph(SUBTITLE)
    doc.add_paragraph(TEAM)
    doc.add_paragraph("Companion files: ACTIVE_WORKSPACE_BOOKLET.md, CHANGE_MAP.md, INTEGRATION_GUIDE.md, RELEASE_CHECKLIST.md, and CHANGE_REQUEST_TEMPLATE.md.")
    for heading, body in SECTIONS:
        doc.add_heading(heading, level=1)
        doc.add_paragraph(body)
    doc.add_heading("Local Portal Screenshot", level=1)
    if SCREENSHOT.exists():
        doc.add_picture(str(SCREENSHOT), width=Inches(6.3))
    doc.add_paragraph("Prepared for the LearnSphere AI internship project.")
    for style in doc.styles:
        if style.name == "Normal":
            style.font.name = "Arial"
            style.font.size = Pt(10)
    doc.save(OUTPUT / "LearnSphere_AI_Active_Workspace_Booklet.docx")


def add_pdf():
    styles = getSampleStyleSheet()
    story = [Paragraph(TITLE, styles["Title"]), Paragraph(SUBTITLE, styles["Heading2"]), Paragraph(TEAM, styles["BodyText"]), Spacer(1, 12)]
    for heading, body in SECTIONS:
        story.extend([Paragraph(heading, styles["Heading2"]), Paragraph(body, styles["BodyText"]), Spacer(1, 8)])
    SimpleDocTemplate(str(OUTPUT / "LearnSphere_AI_Active_Workspace_Booklet.pdf"), pagesize=A4, rightMargin=42, leftMargin=42, topMargin=42, bottomMargin=42).build(story)


def add_slide(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.text = body
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = PptPt(20)


def add_presentation():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "LearnSphere AI\nActive Workspace and Deployment System"
    title.placeholders[1].text = f"{SUBTITLE}\n{TEAM}"
    slides = [
        ("Why an Active Workspace?", "A separate clone protects the baseline project while the team develops, tests, documents, and releases improvements."),
        ("Independent Localhost", "The active workspace runs on http://localhost:5051. It uses its own virtual environment and database, separate from the original project server."),
        ("Readable Page Structure", "Overview, planner, library, practice, and team pages each have their own HTML file. Shared behavior and styling remain in dedicated asset files."),
        ("Integration Architecture", "Browser -> frontend -> Flask API -> SQLite\nGitHub Pages -> runtime config -> Render API\nGit push -> Actions -> GitHub Pages"),
        ("Safe Change Workflow", "Pull approved work. Create a feature branch. Run localhost. Edit mapped files. Test. Review diff. Commit. Push. Review and deploy."),
        ("Team and Internship Identity", "LearnSphere AI Team\nVaishnavi Mali | Nagweni Kumbhar | Aishwarya Ekanath\nChatake Innoworks Pvt. Ltd. | MindForgeAI Internship 2026"),
    ]
    for slide_title, body in slides:
        add_slide(prs, slide_title, body)
    demo = prs.slides.add_slide(prs.slide_layouts[5])
    demo.shapes.title.text = "Live Active Workspace Portal"
    if SCREENSHOT.exists():
        demo.shapes.add_picture(str(SCREENSHOT), PptInches(0.7), PptInches(1.15), width=PptInches(11.9))
    add_slide(prs, "Conclusion", "The active workspace provides a clear path for future development: understandable source files, documented commands, repeatable testing, controlled Git updates, and deployment readiness.")
    prs.save(OUTPUT / "LearnSphere_AI_Active_Workspace_Presentation.pptx")


if __name__ == "__main__":
    add_docx()
    add_pdf()
    add_presentation()
    print(f"Created deliverables in {OUTPUT}")
